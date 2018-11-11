#!/usr/bin/env python
# -*- coding: utf-8 -*-


import os
import tempfile
import argparse
from subprocess import call

from pdf2image import convert_from_path
from pptx import Presentation
from gtts import gTTS


__author__ = ['chaonan99']


## Sometimes ffmpeg is avconv
# FFMPEG_NAME = 'ffmpeg'
FFMPEG_NAME = 'avconv'


def ppt_presenter(pptx_path, pdf_path, output_path):
    with tempfile.TemporaryDirectory() as temp_path:
        images_from_path = convert_from_path(pdf_path)
        prs = Presentation(pptx_path)
        assert len(images_from_path) == len(prs.slides)
        for i, (slide, image) in enumerate(zip(prs.slides, images_from_path)):
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                tts = gTTS(text=notes, lang='en')
                image_path = os.path.join(temp_path, 'frame_{}.jpg'.format(i))
                audio_path = os.path.join(temp_path, 'frame_{}.mp3'.format(i))

                image.save(image_path)
                tts.save(audio_path)

                ffmpeg_call(image_path, audio_path, temp_path, i)

        video_list = [os.path.join(temp_path, 'frame_{}.ts'.format(i)) \
                      for i in range(len(images_from_path))]
        video_list_str = 'concat:' + '|'.join(video_list)
        ffmpeg_concat(video_list_str, output_path)


def ffmpeg_call(image_path, audio_path, temp_path, i):
    out_path_mp4 = os.path.join(temp_path, 'frame_{}.mp4'.format(i))
    out_path_ts = os.path.join(temp_path, 'frame_{}.ts'.format(i))
    call([FFMPEG_NAME, '-loop', '1', '-y', '-i', image_path, '-i', audio_path,
          '-c:v', 'libx264', '-tune', 'stillimage', '-c:a', 'libfdk_aac',
          '-b:a', '192k', '-pix_fmt', 'yuv420p', '-shortest', out_path_mp4])
    call([FFMPEG_NAME, '-y', '-i', out_path_mp4, '-c', 'copy',
          '-bsf:v', 'h264_mp4toannexb', '-f', 'mpegts', out_path_ts])


def ffmpeg_concat(video_list_str, out_path):
    call([FFMPEG_NAME, '-y', '-f', 'mpegts', '-i', '{}'.format(video_list_str),
          '-c', 'copy', '-bsf:a', 'aac_adtstoasc', out_path])


def main():
    parser = argparse.ArgumentParser(description='PPT Presenter help.')
    parser.add_argument('--pptx', help='input pptx path')
    parser.add_argument('--pdf', help='input pdf path')
    parser.add_argument('-o', '--output', help='output path')
    args = parser.parse_args()
    ppt_presenter(args.pptx, args.pdf, args.output)


if __name__ == '__main__':
    main()